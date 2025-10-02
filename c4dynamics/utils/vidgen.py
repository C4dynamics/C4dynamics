# -*- coding: utf-8 -*-
# pragma: no cover

import sys, os, re
# print(sys.version)
import natsort

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from lxml import etree
import tkinter as tk
from tkinter import filedialog






edit_existing = True # False # 
imfol = 'cars2_short/accsize2model' # 'cars2_short/velmodel' # 
delimages = False








# Path to the folder containing images
folder_path = os.path.join('D:/', 'Dropbox', 'c4dynamics', 'examples', '_out', imfol)
# print(folder_path)
# Create a PowerPoint presentation object

if edit_existing: 
    root = tk.Tk()
    root.withdraw()  # Hide the root window

    # Bring the root window to the front
    root.lift()
    root.attributes("-topmost", True)
    root.after_idle(root.attributes, '-topmost', False)


    # Open a file dialog to select the PowerPoint file
    pptfile = filedialog.askopenfilename(
        title="Select PowerPoint File",
        filetypes=[("PowerPoint files", "*.pptx")]
    )

else:
  pptfile = None

prs = Presentation(pptfile)
prs.slide_width = Inches(20)
prs.slide_height = Inches(11.25)


namespace = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}
advance_after = 33 # msec 
duration = 10 # msec 


pattern = re.compile(r'\d+.*\.png$')


filelist = natsort.natsorted(os.listdir(folder_path)) 
# imgfiles = [f for f in dirfiles if f.lower().endswith(('.png', '.jpg', '.jpeg', 'bmp', '.tiff'))] 


# Loop through each file in the folder
cnt = 0
for file_name in filelist:
  if pattern.match(file_name):
    print(file_name)
    # if cnt > 10: break

    if edit_existing:


      slide = prs.slides[cnt]

      shapes_to_remove = []

      for shape in slide.shapes: 
        # if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        #     print('  - AutoShape.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        #     print('  - Chart.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
        #     print('  - Freeform.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        #     print('  - Group.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #     print('  - Picture.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        #     print('  - Placeholder.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        #     print('  - Table.')
        # elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        #     print('  - Text Box.')
        # else:
        #     print('  - This is another type of shape.')

        # print(shape.shape_type.value)
        if shape.shape_type == 13:
          shapes_to_remove.append(shape)
          



      for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
        #   if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        #       print('  - AutoShape removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        #       print('  - Chart removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
        #       print('  - Freeform removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        #       print('  - Group removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #       print('  - Picture removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        #       print('  - Placeholder removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        #       print('  - Table removed.')
        #   elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        #       print('  - Text Box removed.')
        #   else:
        #       print('  - a shape removed')
        # print(shape.shape_type.value, end = '')
        # print('removed')
    
    else: 
      # Add a slide
      slide = prs.slides.add_slide(prs.slide_layouts[6])


    # Define the path to the image file
    img_path = os.path.join(folder_path, file_name)
    # Add the image to the slide
    picture = slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    # slide.shapes._spTree.remove(picture._element)  # Temporarily remove the picture element
    slide.shapes._spTree.insert(2, picture._element)  # Re-insert it at the back

    # set_slide_timing(slide, duration=0.01, advance_after=0.03)
    # Namespaces required for the XML elements
    
    # Access the slide's XML
    slide_xml = slide._element

    # Create timing and transition elements
    timing = etree.SubElement(slide_xml, '{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    
    # Duration of the transition (in milliseconds)
    transition = etree.SubElement(slide_xml, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    transition.set('dur', str(int(duration * 1000)))

    # Set advance timing (in milliseconds)
    advance = etree.SubElement(timing, '{http://schemas.openxmlformats.org/presentationml/2006/main}advAfter')
    advance.text = str(int(advance_after * 1000))

    cnt += 1

# Save the presentation
outfile = pptfile if pptfile else os.path.join(folder_path, os.path.basename(folder_path) + '.pptx') 
try:
  prs.save(outfile)
  print('saved ' + outfile)
except PermissionError:
  prs.save(outfile[:-5] + '_v1' + '.pptx')
  print('saved ' + outfile[:-5] + '_v1' + '.pptx')


if delimages:
  for file_name in filelist:
    if pattern.match(file_name):
        os.remove(os.path.join(folder_path, file_name))

