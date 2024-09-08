from pptx import Presentation
from pptx.util import Inches
import argparse
import natsort
import os, re


def slides_gen(fol): 
  # Path to the folder containing images
  # folder_path = os.path.join('D:/', 'Dropbox', 'c4dynamics', 'examples', 'out', 'cars2_short')
  # print(folder_path)
  # Create a PowerPoint presentation object
  prs = Presentation()
  # TODO take the size from the images: 
  prs.slide_width = Inches(20)
  prs.slide_height = Inches(11.25)

  pattern = re.compile(r'\d+.*\.png$')


  filelist = natsort.natsorted(os.listdir(fol)) 
  # imgfiles = [f for f in dirfiles if f.lower().endswith(('.png', '.jpg', '.jpeg', 'bmp', '.tiff'))] 

  # Loop through each file in the folder
  for file_name in filelist:
    if pattern.match(file_name):
      print(file_name)
      # Add a slide
      slide = prs.slides.add_slide(prs.slide_layouts[6])
      # Define the path to the image file
      img_path = os.path.join(fol, file_name)
      # Add the image to the slide
      slide.shapes.add_picture(img_path, Inches(0), Inches(0), width = prs.slide_width, height = prs.slide_height)

  # Save the presentation
  prs.save(os.path.join(fol, 'vidslides.pptx'))





if __name__ == '__main__': 
  
  # Create an ArgumentParser object
  parser = argparse.ArgumentParser()

  # Add arguments to the parser
  parser.add_argument('--fol', required = True, help = 'images folder')
  parser.add_argument('--debug', action = 'store_true', default = False, help = 'whether to run in debug mode')

  # Parse the command-line arguments
  args = parser.parse_args()
  # Convert parsed arguments to a dictionary
  args_dict = vars(args)


  if args.debug: 
    input(f'Run python debugger using process id. \nSelect the pyenv process. \nPress to continue and wait')
  args_dict.pop('debug')

  slides_gen(**args_dict)
  


